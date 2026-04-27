import os
import fitz # PyMuPDF
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
import zipfile

def convert(input_path, target_format, temp_dir):
    filename_without_ext = os.path.splitext(os.path.basename(input_path))[0]
    
    if target_format == 'docx':
        output_path = os.path.join(temp_dir, f"{filename_without_ext}.docx")
        cv = Converter(input_path)
        cv.convert(output_path)
        cv.close()
        return output_path
        
    elif target_format in ['png', 'jpg']:
        # Extract pages as images
        doc = fitz.open(input_path)
        images_dir = os.path.join(temp_dir, "images")
        os.makedirs(images_dir, exist_ok=True)
        
        image_paths = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_path = os.path.join(images_dir, f"page_{page_num + 1}.{target_format}")
            pix.save(img_path)
            image_paths.append(img_path)
            
        doc.close()
        
        if len(image_paths) == 1:
            return image_paths[0]
            
        # Zip multiple images
        zip_path = os.path.join(temp_dir, f"{filename_without_ext}_images.zip")
        with zipfile.ZipFile(zip_path, 'w') as zf:
            for img_path in image_paths:
                zf.write(img_path, os.path.basename(img_path))
                
        return zip_path
        
    elif target_format == 'pptx':
        doc = fitz.open(input_path)
        prs = Presentation()
        # Default slide size
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        blank_slide_layout = prs.slide_layouts[6]
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            
            img_path = os.path.join(temp_dir, f"temp_{page_num}.png")
            pix.save(img_path)
            
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
        doc.close()
        output_path = os.path.join(temp_dir, f"{filename_without_ext}.pptx")
        prs.save(output_path)
        return output_path
        
    else:
        raise NotImplementedError(f"PDF conversion to {target_format} not supported.")
